from pathlib import Path
import csv
import pdfplumber
import config
from config import get_logger

logger = get_logger(__name__)


def extract_text(file_path: Path) -> str:
    """从支持的文件中提取纯文本，控制长度。"""
    suffix = file_path.suffix.lower()

    # 大文件截断提示
    size_hint = ""
    try:
        file_size = file_path.stat().st_size
        if file_size > config.MAX_INGEST_SIZE_BYTES:
            size_hint = f"[文件大小 {file_size / 1024 / 1024:.1f}MB，仅提取前 {config.MAX_TEXT_LENGTH} 字符]\n"
    except Exception:
        pass

    if suffix in {".txt", ".md"}:
        try:
            content = file_path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            content = file_path.read_text(encoding="gbk", errors="ignore")
        content = content.lstrip("\ufeff")
    elif suffix == ".pdf":
        content = _extract_pdf_text(file_path)
    elif suffix == ".xlsx":
        content = _extract_xlsx_text(file_path)
    elif suffix == ".docx":
        content = _extract_docx_text(file_path)
    elif suffix == ".pptx":
        content = _extract_pptx_text(file_path)
    elif suffix == ".csv":
        content = _extract_csv_text(file_path)
    elif suffix == ".doc":
        content = _extract_doc_text(file_path)
    elif suffix == ".ppt":
        content = _extract_ppt_text(file_path)
    else:
        content = ""

    full = (size_hint + content)[: config.MAX_TEXT_LENGTH]
    return full.strip()


def _extract_pdf_text(file_path: Path) -> str:
    """用 pdfplumber 提取 PDF 全部页面的文字。"""
    texts = []
    try:
        with pdfplumber.open(str(file_path)) as pdf:
            for page in pdf.pages:
                txt = page.extract_text()
                if txt:
                    texts.append(txt)
    except Exception as e:
        return f"[PDF extraction failed: {e}]"
    return "\n".join(texts)


def _extract_xlsx_text(file_path: Path) -> str:
    """用 openpyxl 提取 Excel 各工作表的文本内容（带行数截断）。"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(file_path), data_only=True, read_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                row_count += 1
                if row_count > config.MAX_EXCEL_ROWS:
                    rows.append(f"[Sheet {sheet.title} 超过 {config.MAX_EXCEL_ROWS} 行，已截断]")
                    break
                # 将每行非空单元格拼接成字符串
                row_text = "\t".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                parts.append(f"## Sheet: {sheet.title}\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)
    except Exception as e:
        return f"[XLSX extraction failed: {e}]"


def _extract_docx_text(file_path: Path) -> str:
    """用 python-docx 提取 Word 文档段落文本。"""
    try:
        import docx
        document = docx.Document(str(file_path))
        paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
        return "\n".join(paragraphs)
    except Exception as e:
        return f"[DOCX extraction failed: {e}]"


def _extract_pptx_text(file_path: Path) -> str:
    """用 python-pptx 提取 PPT 所有 slide 的 shape 文本。"""
    try:
        import pptx
        prs = pptx.Presentation(str(file_path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            shapes_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    shapes_text.append(shape.text.strip())
            if shapes_text:
                slides_text.append(f"## Slide {i}\n" + "\n".join(shapes_text))
        return "\n\n".join(slides_text)
    except Exception as e:
        return f"[PPTX extraction failed: {e}]"


def _extract_csv_text(file_path: Path) -> str:
    """用标准库 csv 读取 CSV 文本，限制前 500 行。"""
    try:
        rows = []
        with open(file_path, "r", encoding="utf-8", newline="") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= config.MAX_EXCEL_ROWS:
                    rows.append(f"[CSV 超过 {config.MAX_EXCEL_ROWS} 行，已截断]")
                    break
                row_text = "\t".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    rows.append(row_text)
        return "\n".join(rows)
    except UnicodeDecodeError:
        # 尝试 GBK
        try:
            rows = []
            with open(file_path, "r", encoding="gbk", newline="") as f:
                reader = csv.reader(f)
                for i, row in enumerate(reader):
                    if i >= config.MAX_EXCEL_ROWS:
                        rows.append(f"[CSV 超过 {config.MAX_EXCEL_ROWS} 行，已截断]")
                        break
                    row_text = "\t".join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        rows.append(row_text)
            return "\n".join(rows)
        except Exception as e:
            return f"[CSV extraction failed: {e}]"
    except Exception as e:
        return f"[CSV extraction failed: {e}]"


def _extract_doc_text(file_path: Path) -> str:
    """Windows 下优先尝试 pywin32 COM 接口读取 .doc，否则返回友好提示。"""
    try:
        import win32com.client as win32
        word = win32.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(str(file_path.resolve()))
        text = doc.Range().Text
        doc.Close(False)
        word.Quit()
        return text
    except Exception as e:
        logger.warning("DOC COM extraction failed for %s: %s", file_path, e)
        return "[DOC extraction failed: Microsoft Office not available or file is corrupted. Please convert to DOCX.]"


def _extract_ppt_text(file_path: Path) -> str:
    """Windows 下优先尝试 pywin32 COM 接口读取 .ppt，否则返回友好提示。"""
    try:
        import win32com.client as win32
        ppt = win32.Dispatch("PowerPoint.Application")
        ppt.Visible = False
        presentation = ppt.Presentations.Open(str(file_path.resolve()), WithWindow=False)
        slides_text = []
        for i, slide in enumerate(presentation.Slides, 1):
            shapes_text = []
            for shape in slide.Shapes:
                if hasattr(shape, "TextFrame"):
                    text_frame = shape.TextFrame
                    if hasattr(text_frame, "TextRange"):
                        txt = text_frame.TextRange.Text.strip()
                        if txt:
                            shapes_text.append(txt)
            if shapes_text:
                slides_text.append(f"## Slide {i}\n" + "\n".join(shapes_text))
        presentation.Close()
        ppt.Quit()
        return "\n\n".join(slides_text)
    except Exception as e:
        logger.warning("PPT COM extraction failed for %s: %s", file_path, e)
        return "[PPT extraction failed: Microsoft Office not available or file is corrupted. Please convert to PPTX.]"
